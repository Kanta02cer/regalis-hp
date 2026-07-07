# -*- coding: utf-8 -*-
"""Trillion Bank — 会社紹介ワンパゲャー（A4縦・1枚）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLACK=RGBColor(0x08,0x28,0x58); INK=RGBColor(0x1C,0x1C,0x1E); MUTED=RGBColor(0x76,0x74,0x70)
GOLD=RGBColor(0xC0,0x90,0x30); GOLDD=RGBColor(0x94,0x6C,0x24); WHITE=RGBColor(0xFF,0xFF,0xFF)
SOFT=RGBColor(0xF3,0xF5,0xF9); SOFT2=RGBColor(0xE7,0xEC,0xF3); LINE=RGBColor(0xD7,0xDE,0xE8)
JP="游ゴシック"; SER="Georgia"
LOGO_DIR="/home/user/regalis-hp/investor-materials/"
SYMD=LOGO_DIR+"trillion-symbol-dark.png"
def pic(path,l,t,w=None,h=None):
    kw={}
    if w is not None: kw['width']=Inches(w)
    if h is not None: kw['height']=Inches(h)
    return s.shapes.add_picture(path,Inches(l),Inches(t),**kw)

prs=Presentation()
prs.slide_width=Inches(8.27); prs.slide_height=Inches(11.69)
s=prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE

def rect(l,t,w,h,fill=None,line=None,lw=0.75,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    return sp

def txt(l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=2,ls=1.1):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0)
        try: p.line_spacing=ls
        except: pass
        for (text,size,color,bold,font) in para:
            r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.color.rgb=color
            r.font.bold=bold; r.font.name=font
            rPr=r._r.get_or_add_rPr(); ea=rPr.find(qn('a:ea'))
            if ea is None: ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
            ea.set('typeface',JP)
    return tb
def P(*r): return list(r)

M=0.55  # margin
W=8.27-2*M

# ── ヘッダー（黒帯） ──
rect(0,0,8.27,1.65,fill=BLACK)
rect(0,1.65,8.27,0.06,fill=GOLD)
pic(SYMD, 6.95, 0.28, h=1.12)
txt(M,0.3,W,0.4,[P(("TRILLION BANK",13,GOLD,True,SER))])
txt(M,0.66,W,0.6,[P(("AIに選ばれる企業を、設計する。",23,WHITE,True,JP))])
txt(M,1.22,W,0.35,[P(("AI検索に特化したマーケティング & デジタルPRグループ",11,RGBColor(0xCF,0xC4,0xAC),False,JP))])

# ── リード文 ──
y=1.95
txt(M,y,W,0.7,[P(("検索する人の50%以上がAI検索を使う時代、日本語圏サイトの90%以上はまだ未対策。",11,INK,False,JP)),
               P(("私たちは、AIに引用される仕組みを「設計」から自動でつくり、収益に直結させます。",11,INK,False,JP))],ls=1.3)
y=2.75

# ── HackⅡ トリオ ──
txt(M,y,W,0.3,[P(("主力プロダクト：全自動AI検索最適化インフラ HackⅡ（ハックツ）",12,GOLDD,True,JP))]); y+=0.42
tri=[("ハカル","AI引用モニタリング。5モデル横断で引用シェアを可視化。"),
     ("ツクル","情報供給インフラ（特許出願中）。llms.txtを自動生成・更新。"),
     ("ツナグ","MQL顧客アプローチ。AI経由リードの成約率は4.4倍。")]
cw=(W-0.3)/3
for i,(h,b) in enumerate(tri):
    lx=M+i*(cw+0.15)
    rect(lx,y,cw,1.15,fill=SOFT,line=LINE,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(lx,y,cw,0.34,fill=BLACK)
    txt(lx,y+0.03,cw,0.3,[P((h,12,GOLD,True,JP))],align=PP_ALIGN.CENTER)
    txt(lx+0.14,y+0.42,cw-0.28,0.7,[P((b,8.8,INK,False,JP))],ls=1.16)
y+=1.4

# ── コア事業＋料金 ──
txt(M,y,W,0.3,[P(("Core AI 4事業",12,GOLDD,True,JP))]); y+=0.4
biz=[("SEO・AIOメディア運営","★メイン／月額¥98,000〜（税別）"),
     ("AI検索ツール HackⅡ","¥9,800〜/月・タグ1行で即日導入"),
     ("AI検索クローラー収益化","アドセンス型（AI Exclusive™）"),
     ("AI・DX戦略コンサル","設計から始める一気通貫支援")]
cw2=(W-0.15)/2
for i,(h,b) in enumerate(biz):
    lx=M+(i%2)*(cw2+0.15); ty=y+(i//2)*0.72
    rect(lx,ty,cw2,0.62,fill=WHITE,line=LINE,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(lx,ty,0.07,0.62,fill=GOLD)
    txt(lx+0.2,ty+0.08,cw2-0.3,0.3,[P((h,10.5,INK,True,JP))])
    txt(lx+0.2,ty+0.34,cw2-0.3,0.25,[P((b,8.5,MUTED,False,JP))])
y+=0.72*2+0.15

# ── 実績メトリクス ──
rect(M,y,W,1.1,fill=SOFT,line=LINE,lw=1.0)
mets=[("134件+","法人マーケ実績"),("4.4倍","AI経由成約率"),("+26.9pt","AICS改善"),("8名","外部顧問"),("特許","出願中")]
mw=W/len(mets)
for i,(v,l) in enumerate(mets):
    mx=M+i*mw
    txt(mx,y+0.18,mw,0.45,[P((v,17,GOLDD,True,SER))],align=PP_ALIGN.CENTER)
    txt(mx,y+0.68,mw,0.3,[P((l,8.5,INK,False,JP))],align=PP_ALIGN.CENTER)
    if i>0: rect(mx,y+0.2,0.008,0.7,fill=LINE)
y+=1.28

# ── 代表 ──
txt(M,y,W,0.3,[P(("代表取締役CEO",12,GOLDD,True,JP))]); y+=0.4
rect(M,y,W,1.25,fill=WHITE,line=LINE,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(M+0.25,y+0.15,4,0.35,[P(("井上 幹太　",15,INK,True,JP),("Kanta Inoue",10,GOLDD,False,SER))])
txt(M+0.25,y+0.5,W-0.5,0.7,
    [P(("12歳から12年間の不登校を経て独学で起業した「機会建築家」。",9.5,INK,False,JP)),
     P(("令和の虎 累計1,600万円・2連続完全ALL／ソフトバンクアカデミア17期／J-StarX第1期／ZEN大学1期特別奨学生。",9.5,MUTED,False,JP))],ls=1.25)
y+=1.45

# ── 会社概要フッター ──
rect(M,y,W,1.15,fill=BLACK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(M+0.25,y+0.14,W-0.5,0.3,[P(("会社概要 / Contact",11,GOLD,True,JP))])
info=[("会社名：Trillion Bank"),
      ("所在地：東京都千代田区麹町6丁目2-1 麹町サイトビル6階"),
      ("Web：regalis-order-suits.com　｜　無料AI引用診断（30分・費用なし）受付中")]
yy=y+0.45
for t in info:
    txt(M+0.25,yy,W-0.5,0.24,[P((t,9.3,RGBColor(0xE6,0xE1,0xD6),False,JP))]); yy+=0.23

out="/home/user/regalis-hp/investor-materials/TrillionBank_会社紹介ワンパゲャー.pptx"
prs.save(out); print("saved:",out)
