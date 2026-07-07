# -*- coding: utf-8 -*-
"""
Regalis Japan Group — 投資家向けピッチ資料（事業計画書）ビルダー
スライド設計ナレッジ（ワンスライド・ワンメッセージ／主張文タイトル／3色ルール）と
財務会計ナレッジ（前提分離・数式駆動・単位/出所明記・想定値の明示）に準拠。

※ 財務数値はすべて「前提を明示した想定モデル（ベースシナリオ）」であり、
   実績値ではありません。各前提はAssumptionスライドに集約しています。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── ブランド・カラーパレット（3色ルール：黒＋金＋淡色ベース） ──────────
BLACK  = RGBColor(0x10, 0x10, 0x12)   # カバー背景
INK    = RGBColor(0x1C, 0x1C, 0x1E)   # 本文
MUTED  = RGBColor(0x76, 0x74, 0x70)   # サブテキスト
GOLD   = RGBColor(0xC5, 0xA0, 0x59)   # ブランドゴールド（アクセント）
GOLDD  = RGBColor(0x9C, 0x7C, 0x3C)   # 濃ゴールド
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SOFT   = RGBColor(0xF7, 0xF3, 0xEB)   # 淡ゴールドパネル
SOFT2  = RGBColor(0xEF, 0xEA, 0xDE)
LINE   = RGBColor(0xE2, 0xDA, 0xC9)   # 罫線
GREEN  = RGBColor(0x3B, 0x7A, 0x57)   # プラス
RED    = RGBColor(0xB0, 0x44, 0x2F)   # マイナス／警告

JP   = "游ゴシック"
JPB  = "游ゴシック"
SER  = "Georgia"           # 英字セリフのアクセント（表紙・数値）

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

# ── ヘルパー ──────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, l, t, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.06, wrap=True):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold,font,italic)"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        try: p.line_spacing = line_spacing
        except Exception: pass
        for (text, size, color, bold, font, *rest) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font
            if rest and rest[0]: r.font.italic = True
    return tb

def set_ea(tb, font=JP):
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
            ea.set('typeface', font)

def P(*runs):  # single paragraph shortcut
    return list(runs)

# footer / page number on light slides
PAGE = {"n": 0}
def footer(s, light=True):
    PAGE["n"] += 1
    c = MUTED if light else RGBColor(0x8A,0x86,0x7E)
    t = txt(s, 0.55, 7.02, 6, 0.3,
            [P(("Regalis Japan Group  ｜  事業計画書（Confidential）", 8, c, False, JP))])
    set_ea(t)
    t2 = txt(s, 11.5, 7.02, 1.3, 0.3,
             [P((f"{PAGE['n']:02d}", 9, c, False, SER))], align=PP_ALIGN.RIGHT)

def kicker(s, label, color=GOLDD):
    t = txt(s, 0.6, 0.5, 8, 0.3,
            [P((label, 10.5, color, True, JP))])
    set_ea(t)
    rect(s, 0.62, 0.86, 0.55, 0.045, fill=GOLD)

def title(s, main, sub=None, y=0.98):
    t = txt(s, 0.6, y, 12.1, 1.0,
            [P((main, 25, INK, True, JPB))], line_spacing=1.05)
    set_ea(t, JPB)
    if sub:
        t2 = txt(s, 0.6, y+0.72, 12.1, 0.5, [P((sub, 12.5, MUTED, False, JP))])
        set_ea(t2)

def chip(s, l, t, w, h, text, fill, fg, size=11, bold=True):
    r = rect(s, l, t, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    r.adjustments[0] = 0.18
    tf = r.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rn = p.add_run(); rn.text = text
    rn.font.size = Pt(size); rn.font.color.rgb = fg; rn.font.bold = bold; rn.font.name = JP
    set_ea(r, JP)
    return r

def card(s, l, t, w, h, fill=WHITE, line=LINE):
    return rect(s, l, t, w, h, fill=fill, line=line, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)

# ══════════════════════════════════════════════════════════
# 01 表紙
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, BLACK)
rect(s, 0, 0, 13.333, 0.16, fill=GOLD)
t = txt(s, 0.9, 1.5, 11, 0.5, [P(("REGALIS JAPAN GROUP", 15, GOLD, True, SER))])
set_ea(t)
t = txt(s, 0.86, 2.15, 11.6, 1.6,
        [P(("AIに選ばれる企業を、", 40, WHITE, True, JPB)),
         P(("設計する。", 40, WHITE, True, JPB))], line_spacing=1.08)
set_ea(t, JPB)
t = txt(s, 0.9, 4.35, 11, 0.5,
        [P(("事業計画書 ／ Business Plan 2026–2028", 14, RGBColor(0xCF,0xC4,0xAC), False, JP))])
set_ea(t)
rect(s, 0.9, 5.02, 3.4, 0.02, fill=GOLDD)
t = txt(s, 0.9, 5.25, 11.5, 1.2,
        [P(("日本初の全自動AI検索最適化インフラ「HackⅡ（ハックツ）」で、", 13.5, RGBColor(0xC8,0xC5,0xBE), False, JP)),
         P(("ChatGPT・Perplexity・Geminiに引用される企業をつくる。", 13.5, RGBColor(0xC8,0xC5,0xBE), False, JP))],
        line_spacing=1.3)
set_ea(t)
t = txt(s, 0.9, 6.65, 11.5, 0.5,
        [P(("Regalis Japan Group株式会社　｜　東京都千代田区麹町　｜　代表取締役 井上 幹太", 10.5, RGBColor(0x9B,0x93,0x82), False, JP))])
set_ea(t)
t = txt(s, 0.9, 6.98, 11.5, 0.4,
        [P(("本資料は投資家向けの事業計画書です。財務数値は前提を明示した想定モデルであり実績値ではありません。", 8, RGBColor(0x77,0x71,0x64), False, JP))])
set_ea(t)

# ══════════════════════════════════════════════════════════
# 02 エグゼクティブサマリー
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "EXECUTIVE SUMMARY")
title(s, "検索がAIに移行する転換点で、「AIに引用される仕組み」を自動化する")
# 3本柱
cols = [
    ("01  巨大な地殻変動", "検索する人の50%以上がAI検索を利用する一方、日本語圏サイトの90%以上がAI検索に未対策。「対策不能」ではなく「手法が確立していない」だけの空白市場。", GOLD),
    ("02  独自プロダクト", "タグ1行で導入でき、AIクローラーを自動検知してLLMOを継続実装する全自動インフラ「HackⅡ」。動的更新エンジンは特許出願中で、模倣困難な技術的堀を形成。", GOLDD),
    ("03  収益に直結", "「計測」で終わらせず、AI引用経由の高単価リード獲得（成約率4.4倍）とデータライセンスまで設計。ストック型SaaS×高粗利で拡張。", GOLD),
]
x = 0.6; w = 3.95; gap = 0.19
for i,(h,b,acc) in enumerate(cols):
    lx = x + i*(w+gap)
    card(s, lx, 2.05, w, 3.05)
    rect(s, lx, 2.05, w, 0.09, fill=acc)
    t = txt(s, lx+0.28, 2.35, w-0.56, 0.6, [P((h, 14.5, INK, True, JPB))]); set_ea(t, JPB)
    t = txt(s, lx+0.28, 3.05, w-0.56, 1.9, [P((b, 11, RGBColor(0x44,0x42,0x40), False, JP))], line_spacing=1.28); set_ea(t)
# ハイライトバー
rect(s, 0.6, 5.45, 12.13, 1.35, fill=SOFT, line=LINE, line_w=1.0)
metrics = [("134件+","法人マーケ実績"),("4.4倍","AI経由MQL成約率"),("+26.9pt","AICSスコア改善"),("8名","外部顧問"),("特許出願中","動的更新エンジン")]
mw = 12.13/len(metrics)
for i,(v,l) in enumerate(metrics):
    mx = 0.6 + i*mw
    t = txt(s, mx, 5.62, mw, 0.6, [P((v, 22, GOLDD, True, SER))], align=PP_ALIGN.CENTER); set_ea(t)
    t = txt(s, mx, 6.28, mw, 0.4, [P((l, 10.5, INK, False, JP))], align=PP_ALIGN.CENTER); set_ea(t)
    if i>0: rect(s, mx, 5.7, 0.008, 0.9, fill=LINE)
footer(s)

# ══════════════════════════════════════════════════════════
# 03 私たちは何者か
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "WHAT WE DO")
title(s, "AI検索に特化したマーケティング & デジタルPRグループ",
      "「見えないブランド価値を、AI時代の定量資産へ」")
rows = [
    ("御社は、AIに選ばれているか。","ChatGPT・Claude・Gemini・Perplexityが「業界で信頼できる会社は？」と聞かれたとき、御社が引用されるかどうか。それが次の10年の集客を左右する。"),
    ("私たちがすること","アナログで曖昧なブランド価値を独自アルゴリズムで定量化し、llms.txt・構造化データとしてAIに供給。AIが最も推薦したくなる状態を継続構築する。"),
    ("設計から始める","「とりあえずツール導入」ではなく、代表の設計思想を軸に、計測→情報供給→商談化までを一気通貫で設計する。"),
]
y = 2.35
for i,(h,b) in enumerate(rows):
    rect(s, 0.6, y, 0.09, 1.35, fill=GOLD)
    t = txt(s, 0.95, y+0.03, 11.5, 0.5, [P((h, 16, INK, True, JPB))]); set_ea(t, JPB)
    t = txt(s, 0.95, y+0.55, 11.5, 0.8, [P((b, 12, RGBColor(0x44,0x42,0x40), False, JP))], line_spacing=1.3); set_ea(t)
    y += 1.55
footer(s)

# ══════════════════════════════════════════════════════════
# 04 市場機会
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "MARKET OPPORTUNITY")
title(s, "検索の主役がAIに移る「今」が、市場を取りに行く唯一のタイミング")
# 左：3つの数字
lefts = [("50%+","検索する人のうちAI検索を利用する割合（2026年前後推定）"),
         ("90%+","日本語圏サイトのAI検索未対策率（Search Engine Journal調査）"),
         ("70%","BtoB購買決定がセールス接触前にデジタルで完結（Gartner調査）")]
y=2.15
for v,l in lefts:
    rect(s,0.6,y,5.7,1.42, fill=SOFT, line=LINE, line_w=1.0)
    t=txt(s,0.85,y+0.16,2.2,1.1,[P((v,30,GOLDD,True,SER))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
    t=txt(s,2.75,y+0.16,3.4,1.15,[P((l,11,INK,False,JP))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.24); set_ea(t)
    y+=1.6
# 右：TAM/SAM/SOM 同心
rx=6.7
card(s,rx,2.05,6.03,4.5)
t=txt(s,rx+0.3,2.25,5.5,0.4,[P(("市場規模（想定レンジ）",13,INK,True,JPB))]); set_ea(t,JPB)
bands=[("TAM","国内デジタルマーケティング市場","約 3兆円規模",GOLD,5.5),
       ("SAM","SEO/コンテンツ＋AI最適化領域","2028年 数千億円へ",GOLDD,4.2),
       ("SOM","AI検索最適化・初期獲得可能市場","数十〜数百億円",INK,2.9)]
yy=2.85
for tag,desc,val,col,bw in bands:
    rect(s,rx+0.3,yy,bw,0.92, fill=SOFT2 if tag!="SOM" else SOFT, line=col, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t=txt(s,rx+0.5,yy+0.08,bw-0.3,0.4,[P((tag+"　",13,col,True,SER),(val,12,INK,True,JP))]); set_ea(t)
    t=txt(s,rx+0.5,yy+0.5,bw-0.3,0.35,[P((desc,9.5,MUTED,False,JP))]); set_ea(t)
    yy+=1.06
t=txt(s,rx+0.3,6.05,5.6,0.4,[P(("※ 各種公開統計を基にした想定レンジ。実数は一次情報での確認が必要。",8,MUTED,False,JP))]); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 05 課題
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "PROBLEM")
title(s, "AI検索時代の勝敗は「引用されるか」で決まるのに、打ち手がない")
probs=[
 ("IR・情報の誤引用リスク","AIが公式資料ではなく競合PRや古い記事を根拠に回答。投資家・取引先に誤情報が拡散する「IR情報逆転現象」。"),
 ("計測ツールで止まる既存策","otterly等の海外ツールは引用状況を後追いで「可視化」するだけ。Google AI Overview偏重で日本語精度に課題。"),
 ("機械可読でない情報資産","IR・製品情報はPDF・スライド中心で、AIが取得できない。価値ある情報ほどAIに届いていない。"),
 ("属人的で継続できない","AIモデルもアルゴリズムも高速に変化。人手での対応は追従できず、一度作って終わりになりがち。"),
]
x=0.6; w=6.0; hgap=0.13; vgap=0.2; h=1.9
for i,(hh,bb) in enumerate(probs):
    lx=x+(i%2)*(w+hgap); ty=2.15+(i//2)*(h+vgap)
    card(s,lx,ty,w,h, fill=WHITE, line=LINE)
    rect(s,lx,ty,0.09,h, fill=RED)
    t=txt(s,lx+0.3,ty+0.22,w-0.55,0.5,[P((hh,14.5,INK,True,JPB))]); set_ea(t,JPB)
    t=txt(s,lx+0.3,ty+0.78,w-0.55,1.0,[P((bb,11,RGBColor(0x44,0x42,0x40),False,JP))],line_spacing=1.26); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 06 ソリューション HackⅡ
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "SOLUTION — HackⅡ")
title(s, "タグ1行で導入、AIへの引用を全自動で作り続けるインフラ「HackⅡ」",
      "「御社の情報を、AIの脳内に直接叩き込む」")
tri=[
 ("ハカル","測る","AI引用モニタリング","5モデル横断で引用シェアを可視化。100クエリを約30分で処理し、競合との相対シェアを数値化。"),
 ("ツクル","創る","情報供給インフラ（特許出願中）","llms.txt・構造化データを自動生成・自動更新。AIに届く情報を継続的に整備し続ける。"),
 ("ツナグ","繋ぐ","MQL顧客アプローチ","引用で終わらせず商談へ。AI経由リードの成約率は通常流入比4.4倍。"),
]
x=0.6; w=3.95; gap=0.19
for i,(jp,mean,role,desc) in enumerate(tri):
    lx=x+i*(w+gap)
    card(s,lx,2.5,w,3.7)
    rect(s,lx,2.5,w,0.7, fill=BLACK)
    t=txt(s,lx,2.62,w,0.5,[P((jp,20,GOLD,True,JPB),("（"+mean+"）",12,WHITE,False,JP))],align=PP_ALIGN.CENTER); set_ea(t,JPB)
    t=txt(s,lx+0.28,3.42,w-0.56,0.6,[P((role,13,INK,True,JPB))],align=PP_ALIGN.CENTER); set_ea(t,JPB)
    rect(s,lx+1.2,4.08,w-2.4,0.02, fill=LINE)
    t=txt(s,lx+0.28,4.25,w-0.56,1.7,[P((desc,11.5,RGBColor(0x44,0x42,0x40),False,JP))],align=PP_ALIGN.CENTER,line_spacing=1.3); set_ea(t)
t=txt(s,0.6,6.45,12,0.4,[P(("料金：Starter ¥9,800〜／Pro ¥29,800〜／Enterprise ¥98,000〜／Platform ¥1,500,000〜（月額・税込）　初期Web開発費 無料（6ヶ月契約前提）",10.5,MUTED,False,JP))]); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 07 技術・堀
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "TECHNOLOGY & MOAT")
title(s, "3層アーキテクチャと「動的更新エンジン」が模倣困難な堀をつくる")
layers=[
 ("Layer 01","データ取得","6層マルチネットワーク。Google検索等の外部API群からAI検索動態データを収集。"),
 ("Layer 02","処理エンジン","独自アルゴリズムで5つの主要AIモデルへ非同期APIクエリを同時発行し最適化。"),
 ("Layer 03","動的更新エンジン（特許出願中）","AIクローラーのトラフィック動態変化を自動検知し、llms.txtを自動再構成。LLMOを完全自動で継続実装する世界初の動的最適化アーキテクチャ。"),
]
y=2.2
for i,(tag,h,b) in enumerate(layers):
    hl = (i==2)
    fill = BLACK if hl else SOFT
    rect(s,0.6,y,12.13,1.32, fill=fill, line=(GOLD if hl else LINE), line_w=1.4 if hl else 1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t=txt(s,0.9,y+0.2,1.9,0.9,[P((tag,15,GOLD if hl else GOLDD,True,SER))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
    t=txt(s,2.9,y+0.16,3.3,1.0,[P((h,14.5,WHITE if hl else INK,True,JPB))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.1); set_ea(t,JPB)
    t=txt(s,6.3,y+0.16,6.15,1.0,[P((b,11.5,(RGBColor(0xD8,0xD3,0xC8) if hl else RGBColor(0x44,0x42,0x40)),False,JP))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.24); set_ea(t)
    y+=1.5
footer(s)

# ══════════════════════════════════════════════════════════
# 08 Why Regalis（差別化の核）
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "WHY REGALIS")
title(s, "「可視化」する競合と違い、私たちは「推薦される状態」をつくり収益化する")
w3=[
 ("資産の定量化","曖昧なブランド価値を独自アルゴリズムで定量化し、AIが推薦したくなる構造化データに変換。"),
 ("中央集権型APIインフラ","SaaSダッシュボードの代理店モデルではなく、タグ1行で即日導入できるAPI駆動の直接ハック。"),
 ("収益に直結","ブランド言及率レポートではなく、AI経由の高単価送客＋データライセンスという新しいキャッシュフロー。"),
]
x=0.6; w=3.95; gap=0.19
for i,(h,b) in enumerate(w3):
    lx=x+i*(w+gap)
    card(s,lx,2.2,w,2.4)
    t=txt(s,lx+0.28,2.45,0.9,0.6,[P((f"0{i+1}",24,GOLD,True,SER))]); set_ea(t)
    t=txt(s,lx+0.28,3.05,w-0.56,0.6,[P((h,14.5,INK,True,JPB))],line_spacing=1.1); set_ea(t,JPB)
    t=txt(s,lx+0.28,3.7,w-0.56,0.9,[P((b,11,RGBColor(0x44,0x42,0x40),False,JP))],line_spacing=1.26); set_ea(t)
# 比較表
th=["比較項目","既存のLLMO/AIO計測ツール","Regalis AIO Intelligence"]
tr=[
 ("提供価値の核","引用状況を後追いで可視化","データを定量化・最適化し推薦される状態を作る"),
 ("対応範囲","AI Overview偏重／日本語精度に課題","日本語特化・ChatGPT/Claude/Perplexity網羅"),
 ("ROI地点","ブランド言及率レポート提出","AI経由の送客＋データ販売収益"),
]
ty=4.85; rowh=0.52; c0=0.6; c1=3.0; c2=5.0; cW=[3.0,4.5,4.63]
xs=[0.6,3.6,8.1]; ws=[3.0,4.5,4.63]
# header
for j,htext in enumerate(th):
    rect(s,xs[j],ty,ws[j],rowh, fill=BLACK, line=WHITE, line_w=0.5)
    t=txt(s,xs[j]+0.12,ty,ws[j]-0.2,rowh,[P((htext,10.5,GOLD if j==2 else WHITE,True,JP))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
for r,row in enumerate(tr):
    yy=ty+rowh*(r+1)
    for j,cell in enumerate(row):
        fill = WHITE if j<2 else SOFT
        rect(s,xs[j],yy,ws[j],rowh, fill=fill, line=LINE, line_w=0.5)
        col = INK if j<2 else GOLDD
        b = (j==2)
        t=txt(s,xs[j]+0.12,yy,ws[j]-0.2,rowh,[P((cell,9.8,col,b,JP))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.05); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 09 ビジネスモデル / ユニットエコノミクス
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "BUSINESS MODEL")
title(s, "ストック型SaaS × 高単価サービス × データ収益の3層マネタイズ")
# 左：収益構造
card(s,0.6,2.15,6.0,4.4)
t=txt(s,0.85,2.35,5.5,0.4,[P(("収益の3層構造",13,INK,True,JPB))]); set_ea(t,JPB)
layers=[
 ("① SaaS（ストック）","HackⅡ 月額課金。Starter〜Platformの階層。低churn・高粗利のリカーリング収益。",GOLD),
 ("② 運用・コンサル（フロー→ストック化）","AIOメディア運営代行 ¥98,000〜/月（6ヶ月契約）、DX・Web開発。",GOLDD),
 ("③ データ・アドセンス型収益","AI検索クローラー収益化（AI Exclusive™）、データライセンス販売。",INK),
]
yy=2.9
for h,b,c in layers:
    rect(s,0.85,yy,0.08,1.05, fill=c)
    t=txt(s,1.05,yy,5.3,0.4,[P((h,12,INK,True,JPB))]); set_ea(t,JPB)
    t=txt(s,1.05,yy+0.4,5.3,0.7,[P((b,10.5,RGBColor(0x44,0x42,0x40),False,JP))],line_spacing=1.22); set_ea(t)
    yy+=1.16
# 右：ユニットエコノミクス（想定）
card(s,6.73,2.15,6.0,4.4, fill=SOFT)
t=txt(s,6.98,2.35,5.5,0.4,[P(("ユニットエコノミクス（想定・HackⅡ SaaS）",13,INK,True,JPB))]); set_ea(t,JPB)
ue=[("平均月次単価 ARPA","¥38,000"),("粗利率","約 82%"),("月次解約率（churn）","2.0%"),
    ("顧客生涯価値 LTV","約 156万円"),("顧客獲得コスト CAC","約 12万円"),("LTV / CAC","約 13x"),("CAC回収期間","約 4ヶ月")]
yy=2.95
for i,(k,v) in enumerate(ue):
    highlight = k in ("LTV / CAC","CAC回収期間")
    t=txt(s,6.98,yy,3.6,0.35,[P((k,11,INK,highlight,JP))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
    t=txt(s,10.3,yy,2.15,0.35,[P((v,13 if highlight else 12,GOLDD if highlight else INK,True,SER))],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
    rect(s,6.98,yy+0.44,5.5,0.008, fill=LINE)
    yy+=0.5
t=txt(s,6.98,6.28,5.5,0.3,[P(("※ 前提に基づく想定値。詳細は「前提条件」スライド参照。",8,MUTED,False,JP))]); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 10 トラクション
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "TRACTION")
title(s, "自社実証で成果を出し、信頼できる顧問陣と実績を積み上げている")
tr=[("134件+","法人向けウェブマーケ実績"),
    ("4.4倍","AI経由MQLの成約率（通常流入比）"),
    ("51→77.9pt","全143記事の平均AICSスコア改善（+26.9pt）"),
    ("8名","外部顧問（住友商事G・SBIグループ・Vector Group等）"),
    ("100クエリ/30分","AI引用モニタリング処理速度"),
    ("5モデル","対応AI（ChatGPT/Perplexity/AI Overview/Claude/Gemini）")]
x=0.6; w=3.95; gap=0.19; h=1.85
for i,(v,l) in enumerate(tr):
    lx=x+(i%3)*(w+gap); ty=2.2+(i//3)*(h+0.22)
    card(s,lx,ty,w,h)
    rect(s,lx,ty,w,0.08, fill=GOLD)
    t=txt(s,lx+0.28,ty+0.28,w-0.56,0.7,[P((v,25,GOLDD,True,SER))]); set_ea(t)
    t=txt(s,lx+0.28,ty+1.02,w-0.56,0.7,[P((l,11,INK,False,JP))],line_spacing=1.22); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 11 事業ポートフォリオ
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "BUSINESS PORTFOLIO")
title(s, "Core AI 4事業を主軸に、展開事業でブランドと人材基盤を広げる")
t=txt(s,0.6,1.95,12,0.35,[P(("■ Core AI 4事業（主軸）",12,GOLDD,True,JP))]); set_ea(t)
core=[("SEO・AIOメディア運営","★メイン事業／月額¥98,000〜"),
      ("AI検索ツール（HackⅡ）","AI引用モニタリング・自動最適化"),
      ("AI検索クローラー収益化","アドセンス型（AI Exclusive™）"),
      ("AI・DX戦略コンサル","設計から始める伴走支援")]
x=0.6; w=2.96; gap=0.13
for i,(h,b) in enumerate(core):
    lx=x+i*(w+gap)
    card(s,lx,2.35,w,1.5, fill=SOFT)
    t=txt(s,lx+0.22,2.52,w-0.4,0.7,[P((h,12,INK,True,JPB))],line_spacing=1.1); set_ea(t,JPB)
    t=txt(s,lx+0.22,3.28,w-0.4,0.5,[P((b,9.8,MUTED,False,JP))],line_spacing=1.15); set_ea(t)
t=txt(s,0.6,4.15,12,0.35,[P(("■ 展開事業",12,GOLDD,True,JP))]); set_ea(t)
exp=[("Web・システム開発",""),("日本学生アンバサダー協会",""),("代表タレント活動・講演",""),("販売促進・営業代行","")]
for i,(h,b) in enumerate(exp):
    lx=x+i*(w+gap)
    card(s,lx,4.55,w,1.2, fill=WHITE)
    t=txt(s,lx+0.22,4.55,w-0.4,1.2,[P((h,11.5,INK,True,JP))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.1); set_ea(t)
t=txt(s,0.6,5.95,12,0.5,[P(("「機会が届いていない場所に、機会を設計し直す」— 展開事業は次世代人材の獲得とブランド接点を生み、Core事業に還流する。",11,MUTED,False,JP))],line_spacing=1.2); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 12 Go-to-Market
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "GO-TO-MARKET")
title(s, "「無料AI引用診断」を入口に、SaaSとデータ収益へ引き上げる導線設計")
steps=[
 ("① 認知・発見","自社メディアのAIO実証・記事資産でAI検索・SNS経由の流入を獲得。代表のタレント活動が信頼を補強。"),
 ("② 無料AI引用診断（30分）","業界内AI引用シェアと機会損失額を可視化。摩擦最小の第一CTAで商談化。"),
 ("③ 導入・運用","HackⅡ／AIOメディア運営で6ヶ月コミット。タグ1行で即日計測、成果を積み上げる。"),
 ("④ 拡張・アップセル","Enterprise（IR/ESG連携）・Platform・データライセンスへ。NRR向上で1顧客あたり収益を拡大。"),
]
y=2.25; h=1.02
for i,(h_,b) in enumerate(steps):
    rect(s,0.6,y,12.13,h, fill=(SOFT if i%2==0 else WHITE), line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    t=txt(s,0.9,y,4.0,h,[P((h_,14,INK,True,JPB))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.1); set_ea(t,JPB)
    t=txt(s,4.9,y,7.6,h,[P((b,11.5,RGBColor(0x44,0x42,0x40),False,JP))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.22); set_ea(t)
    y+=h+0.16
footer(s)

# ══════════════════════════════════════════════════════════
# 13 財務計画（サマリー＋グラフ）
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "FINANCIAL PLAN（想定・ベースシナリオ）")
title(s, "先行投資でFY26は赤字も、FY27に黒字転換しFY28に営業利益率20%超へ")
# チャート：セグメント別売上（積み上げ縦棒）
chart_data = CategoryChartData()
chart_data.categories = ['FY2026', 'FY2027', 'FY2028']
chart_data.add_series('HackⅡ SaaS', (16.0, 73.0, 214.3))
chart_data.add_series('AIOメディア運営', (23.8, 55.4, 103.0))
chart_data.add_series('DX・Web開発', (12.0, 28.0, 48.0))
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,
    Inches(0.6), Inches(2.1), Inches(6.3), Inches(4.3), chart_data)
chart = gframe.chart
chart.has_title = True
chart.chart_title.text_frame.text = "売上高（セグメント別・百万円）"
for run in chart.chart_title.text_frame.paragraphs[0].runs:
    run.font.size=Pt(11); run.font.bold=True; run.font.name=JP; run.font.color.rgb=INK
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size=Pt(9); chart.legend.font.name=JP
plot = chart.plots[0]
series_colors=[GOLD,GOLDD,RGBColor(0xBD,0xB4,0xA2)]
for sidx,ser in enumerate(plot.series):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb=series_colors[sidx]
cat_ax=chart.category_axis; cat_ax.tick_labels.font.size=Pt(10); cat_ax.tick_labels.font.name=JP
val_ax=chart.value_axis; val_ax.tick_labels.font.size=Pt(9)
val_ax.has_major_gridlines=True

# 右：主要指標テーブル
xs=[6.95,9.35,10.45,11.55]; ws=[2.4,1.1,1.1,1.18]; ty=2.25; rowh=0.5
hdr=["百万円 / FY","2026","2027","2028"]
for j,htext in enumerate(hdr):
    rect(s,xs[j],ty,ws[j],rowh, fill=BLACK)
    al = PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
    t=txt(s,xs[j]+0.1,ty,ws[j]-0.2,rowh,[P((htext,10,WHITE,True,SER if j>0 else JP))],anchor=MSO_ANCHOR.MIDDLE,align=al); set_ea(t)
frows=[("売上高",("52","156","365"),INK,False),
       ("売上総利益",("32","103","254"),INK,False),
       ("（粗利率）",("61%","66%","69%"),MUTED,False),
       ("販管費",("42","88","175"),INK,False),
       ("営業利益",("△10","15","79"),None,True),
       ("（営業利益率）",("△20%","10%","22%"),MUTED,False),
       ("従業員数（期末）",("12","18","32"),INK,False)]
for r,(k,vals,col,emp) in enumerate(frows):
    yy=ty+rowh*(r+1)
    rect(s,xs[0],yy,ws[0],rowh, fill=(SOFT if emp else WHITE), line=LINE, line_w=0.5)
    t=txt(s,xs[0]+0.1,yy,ws[0]-0.2,rowh,[P((k,9.8,col or INK,emp,JP))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
    for j,v in enumerate(vals):
        rect(s,xs[j+1],yy,ws[j+1],rowh, fill=(SOFT if emp else WHITE), line=LINE, line_w=0.5)
        vc = col or INK
        if "△" in v: vc = RED
        elif emp: vc = GREEN
        t=txt(s,xs[j+1]+0.05,yy,ws[j+1]-0.12,rowh,[P((v,10.5 if emp else 10,vc,emp,SER))],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
t=txt(s,6.95,6.35,5.8,0.5,[P(("※ 前提を明示した想定モデル（ベースシナリオ）。実績値ではありません。単位：百万円、△＝マイナス。",8,MUTED,False,JP))],line_spacing=1.15); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 14 前提条件（Assumptions）
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "KEY ASSUMPTIONS")
title(s, "財務計画の前提（すべて想定値・変更可能なドライバーとして分離）")
cols=[
 ("収益ドライバー",[
   "HackⅡ有料アカウント（期末）：60→260→620",
   "AIOメディア運営クライアント（期末）：24→52→88",
   "HackⅡ平均単価 ARPA：¥38,000/月",
   "メディア運営単価：¥110,000/月（本体＋オプション）",
   "DX・Web開発：プロジェクト収益 12→28→48百万円",
 ]),
 ("コスト・効率前提",[
   "粗利率：SaaS 82%／メディア 55%／開発 45%",
   "月次解約率（churn）：2.0%、NRR：115%想定",
   "CAC：¥120,000／LTV：約156万円",
   "従業員数：12→18→32名（採用先行）",
   "販管費：人件費・採用・広告・R&Dを含む",
 ]),
]
x=0.6; w=6.0; gap=0.13
for i,(h,items) in enumerate(cols):
    lx=x+i*(w+gap)
    card(s,lx,2.15,w,4.3)
    rect(s,lx,2.15,w,0.6, fill=BLACK)
    t=txt(s,lx+0.25,2.15,w-0.4,0.6,[P((h,13,GOLD,True,JPB))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t,JPB)
    yy=2.95
    for it in items:
        rect(s,lx+0.28,yy+0.13,0.09,0.09, fill=GOLD, shape=MSO_SHAPE.OVAL)
        t=txt(s,lx+0.5,yy,w-0.8,0.62,[P((it,11,INK,False,JP))],line_spacing=1.18); set_ea(t)
        yy+=0.66
t=txt(s,0.6,6.6,12,0.4,[P(("※ 上記は将来の達成を保証するものではなく、事業設計上のベースシナリオ。市場環境・実績に応じて更新する。",8.5,MUTED,False,JP))]); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 15 マイルストーン（詳細ロードマップ）
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "MILESTONES — ROADMAP")
title(s, "基盤構築 → PMF → スケール → 多角化の4フェーズで段階的に拡大")
phases=[
 ("Phase 1","〜2026 H1","基盤構築",GOLD,[
   "HackⅡ 正式リリース","動的更新エンジン 特許出願",
   "llms.txt 5ファイル体制・ページ別最適化","法人マーケ実績 134件+ を基盤化"]),
 ("Phase 2","2026 H2","製品化・PMF",GOLDD,[
   "Starter/Pro/Enterprise 3プラン提供","ハカル即日稼働ダッシュボード",
   "AICSスコア改善実証（51→77.9）","有料アカウント 〜60"]),
 ("Phase 3","2027","スケール",INK,[
   "Platform・API連携／IR・ESG Enterprise獲得","代理店・パートナー網構築",
   "顧客数×3・営業CF黒字化","NRR 115%・アップセル本格化"]),
 ("Phase 4","2028","拡大・多角化",RED,[
   "AI検索クローラー収益化（AI Exclusive™）","データライセンス収益の立ち上げ",
   "営業利益率 20%超","展開事業（学生協会等）の拡張"]),
]
x=0.6; w=2.96; gap=0.13
# timeline bar
rect(s,0.6,2.15,12.13,0.06, fill=LINE)
for i,(ph,period,name,acc,items) in enumerate(phases):
    lx=x+i*(w+gap)
    rect(s,lx+w/2-0.09,2.06,0.18,0.18, fill=acc, shape=MSO_SHAPE.OVAL)
    card(s,lx,2.45,w,3.95)
    rect(s,lx,2.45,w,0.86, fill=acc)
    t=txt(s,lx,2.52,w,0.4,[P((ph,13,WHITE,True,SER))],align=PP_ALIGN.CENTER); set_ea(t)
    t=txt(s,lx,2.9,w,0.36,[P((period+"　"+name,11,WHITE,True,JP))],align=PP_ALIGN.CENTER); set_ea(t)
    yy=3.5
    for it in items:
        rect(s,lx+0.2,yy+0.1,0.08,0.08, fill=acc, shape=MSO_SHAPE.OVAL)
        t=txt(s,lx+0.4,yy,w-0.55,0.7,[P((it,10,INK,False,JP))],line_spacing=1.15); set_ea(t)
        yy+=0.72
footer(s)

# ══════════════════════════════════════════════════════════
# 16 KPIマイルストーン（数値の里程標）
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "KPI MILESTONES（想定）")
title(s, "各フェーズの到達目標をKPIで定義し、達成度で意思決定する")
hdr=["KPI","FY2026","FY2027","FY2028"]
rows=[
 ("売上高（百万円）","52","156","365"),
 ("ARR ラン・レート（百万円）","59","187","399"),
 ("HackⅡ有料アカウント（期末）","60","260","620"),
 ("メディア運営クライアント（期末）","24","52","88"),
 ("営業利益（百万円）","△10","15","79"),
 ("営業利益率","△20%","10%","22%"),
 ("従業員数（期末）","12","18","32"),
]
xs=[0.6,6.1,8.35,10.6]; ws=[5.5,2.25,2.25,2.13]; ty=2.15; rowh=0.55
for j,h in enumerate(hdr):
    rect(s,xs[j],ty,ws[j],rowh, fill=BLACK)
    al=PP_ALIGN.LEFT if j==0 else PP_ALIGN.RIGHT
    t=txt(s,xs[j]+0.15,ty,ws[j]-0.3,rowh,[P((h,11,WHITE if j==0 else GOLD,True,JP if j==0 else SER))],anchor=MSO_ANCHOR.MIDDLE,align=al); set_ea(t)
for r,(k,*vals) in enumerate(rows):
    yy=ty+rowh*(r+1)
    rect(s,xs[0],yy,ws[0],rowh, fill=(SOFT if r%2 else WHITE), line=LINE, line_w=0.5)
    t=txt(s,xs[0]+0.15,yy,ws[0]-0.3,rowh,[P((k,11,INK,False,JP))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
    for j,v in enumerate(vals):
        rect(s,xs[j+1],yy,ws[j+1],rowh, fill=(SOFT if r%2 else WHITE), line=LINE, line_w=0.5)
        vc = RED if "△" in v else INK
        t=txt(s,xs[j+1]+0.1,yy,ws[j+1]-0.25,rowh,[P((v,11.5,vc,True,SER))],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
t=txt(s,0.6,6.55,12,0.4,[P(("※ すべて前提に基づく想定値（ベースシナリオ）。単位・定義は財務計画スライドに準拠。",8.5,MUTED,False,JP))]); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 17 経営陣・体制
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "TEAM")
title(s, "「機会建築家」が率いる経営と、実力ある外部顧問8名の支援体制")
# 左：代表
card(s,0.6,2.15,6.0,4.4)
rect(s,0.6,2.15,6.0,0.9, fill=BLACK)
t=txt(s,0.9,2.28,5.5,0.4,[P(("井上 幹太　",17,WHITE,True,JPB),("Kanta Inoue",11,GOLD,False,SER))]); set_ea(t,JPB)
t=txt(s,0.9,2.68,5.5,0.3,[P(("代表取締役CEO ／ 機会建築家",11,GOLD,False,JP))]); set_ea(t)
bio=["12歳から12年間の不登校を経て独学で起業。","「機会が届いていない場所に、機会を設計し直す」をミッションに持株会社を設立。"]
yy=3.25
for b in bio:
    t=txt(s,0.9,yy,5.4,0.6,[P((b,11,RGBColor(0x44,0x42,0x40),False,JP))],line_spacing=1.24); set_ea(t)
    yy+=0.62
awards=["令和の虎 Tiger Funding 累計1,600万円・2連続完全ALL",
        "JCI JAPAN TOYP2026 ファイナリスト（青年版国民栄誉賞）",
        "ソフトバンクアカデミア17期 修了",
        "J-StarX（経済産業省）第1期／ZEN大学1期特別奨学生"]
yy=4.55
for a in awards:
    rect(s,0.9,yy+0.1,0.09,0.09, fill=GOLD, shape=MSO_SHAPE.OVAL)
    t=txt(s,1.12,yy,5.2,0.45,[P((a,10.3,INK,False,JP))],line_spacing=1.12); set_ea(t)
    yy+=0.48
# 右：顧問・体制
card(s,6.73,2.15,6.0,4.4, fill=SOFT)
t=txt(s,6.98,2.35,5.5,0.4,[P(("外部顧問・支援体制",13,INK,True,JPB))]); set_ea(t,JPB)
advs=[("外部顧問 8名","住友商事グループ・SBIグループ・Vector Group 等の実務家が支援"),
      ("設計思想を持つ専任チーム","戦略設計→AI実装→組織定着まで一気通貫で伴走"),
      ("代表直接対応","初回商談は代表が担当し、意思決定者に直接価値を届ける")]
yy=2.95
for h,b in advs:
    rect(s,6.98,yy,0.08,0.95, fill=GOLD)
    t=txt(s,7.18,yy,5.3,0.4,[P((h,12.5,INK,True,JPB))]); set_ea(t,JPB)
    t=txt(s,7.18,yy+0.42,5.3,0.6,[P((b,10.8,RGBColor(0x44,0x42,0x40),False,JP))],line_spacing=1.22); set_ea(t)
    yy+=1.12
footer(s)

# ══════════════════════════════════════════════════════════
# 18 リスクと対策
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "RISKS & MITIGATION")
title(s, "主要リスクを特定し、事業設計であらかじめ手当てする")
risks=[
 ("AIプラットフォーム依存","主要LLMの仕様変更・引用ロジック変化","動的更新エンジンで自動追従。5モデル分散で単一依存を回避。"),
 ("市場の未成熟・啓蒙コスト","AIO/LLMOの認知が発展途上","自社実証・記事資産・無料診断で市場を教育しながら獲得。"),
 ("競合・大手参入","海外ツール・大手の日本語対応","日本語特化・API直接ハック・特許出願で差別化を維持。"),
 ("採用・組織拡大","専門人材の獲得と定着","展開事業（学生協会等）を人材パイプラインとして活用。"),
 ("コンプライアンス／IR","誤引用・インサイダー情報混入","ISMS準拠設計・開示連動の自動更新・ロールベース権限管理。"),
 ("キャッシュフロー","先行投資期の資金負担","SaaSのストック収益で早期に回収基盤を構築、churn管理を徹底。"),
]
x=0.6; w=6.0; hgap=0.13; h=1.32
for i,(t1,t2,t3) in enumerate(risks):
    lx=x+(i%2)*(w+hgap); ty=2.1+(i//2)*(h+0.16)
    card(s,lx,ty,w,h)
    rect(s,lx,ty,0.09,h, fill=GOLD)
    t=txt(s,lx+0.28,ty+0.14,w-0.5,0.35,[P((t1,12.5,INK,True,JPB))]); set_ea(t,JPB)
    t=txt(s,lx+0.28,ty+0.5,w-0.5,0.35,[P(("リスク：",9.5,RED,True,JP),(t2,9.8,MUTED,False,JP))],line_spacing=1.1); set_ea(t)
    t=txt(s,lx+0.28,ty+0.86,w-0.5,0.4,[P(("対策：",9.5,GREEN,True,JP),(t3,9.8,RGBColor(0x44,0x42,0x40),False,JP))],line_spacing=1.12); set_ea(t)
footer(s)

# ══════════════════════════════════════════════════════════
# 19 ビジョン
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, BLACK)
rect(s,0,0,13.333,0.14, fill=GOLD)
t=txt(s,0.9,1.35,11.5,0.4,[P(("VISION",13,GOLD,True,SER))]); set_ea(t)
t=txt(s,0.86,2.0,11.6,2.0,
      [P(("AI時代の情報インフラを、",34,WHITE,True,JPB)),
       P(("日本から世界の標準に。",34,WHITE,True,JPB))],line_spacing=1.12)
set_ea(t,JPB)
t=txt(s,0.9,4.2,11.5,1.4,
      [P(("見えないブランド価値を定量化し、AIに正しく届ける。",14,RGBColor(0xCF,0xC4,0xAC),False,JP)),
       P(("「御社は、AIに選ばれているか。」に、すべての企業が自信を持って",14,RGBColor(0xCF,0xC4,0xAC),False,JP)),
       P(("YESと答えられる世界を、設計から始めてつくる。",14,RGBColor(0xCF,0xC4,0xAC),False,JP))],line_spacing=1.4)
set_ea(t)
rect(s,0.9,6.0,3.4,0.02, fill=GOLDD)
t=txt(s,0.9,6.25,11.5,0.4,[P(("Regalis Japan Group株式会社",12,GOLD,True,JP))]); set_ea(t)

# ══════════════════════════════════════════════════════════
# 20 Appendix / お問い合わせ
# ══════════════════════════════════════════════════════════
s = slide(); bg(s, WHITE)
kicker(s, "CONTACT / APPENDIX")
title(s, "お問い合わせ・会社概要")
card(s,0.6,2.15,6.0,4.3, fill=SOFT)
t=txt(s,0.9,2.4,5.5,0.4,[P(("会社概要",13,INK,True,JPB))]); set_ea(t,JPB)
info=[("会社名","Regalis Japan Group株式会社（レガリス）"),
      ("代表者","代表取締役 井上 幹太"),
      ("所在地","東京都千代田区麹町6丁目2-1 麹町サイトビル6階"),
      ("事業内容","AI検索最適化（LLMO/AIO）・DX支援・メディア運営"),
      ("主力製品","全自動AI検索最適化インフラ HackⅡ"),
      ("Web","regalis-order-suits.com")]
yy=2.9
for k,v in info:
    t=txt(s,0.9,yy,1.7,0.5,[P((k,10.5,GOLDD,True,JP))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t)
    t=txt(s,2.55,yy,3.9,0.5,[P((v,10.5,INK,False,JP))],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.15); set_ea(t)
    rect(s,0.9,yy+0.52,5.4,0.008, fill=LINE)
    yy+=0.57
card(s,6.73,2.15,6.0,4.3)
rect(s,6.73,2.15,6.0,0.9, fill=BLACK)
t=txt(s,6.98,2.28,5.5,0.6,[P(("無料AI引用診断（30分）",16,GOLD,True,JPB))],anchor=MSO_ANCHOR.MIDDLE); set_ea(t,JPB)
t=txt(s,6.98,3.25,5.5,1.0,
      [P(("御社の業界内AI引用シェアと機会損失額を可視化します。",11.5,INK,False,JP)),
       P(("費用・義務は一切ありません。初回は代表が直接対応します。",11.5,INK,False,JP))],line_spacing=1.3); set_ea(t)
chip(s,6.98,4.5,5.5,0.7,"regalis-order-suits.com/contact/", GOLD, BLACK, size=12)
t=txt(s,6.98,5.5,5.5,0.8,[P(("本資料はConfidentialです。無断転載・二次配布を禁じます。",9.5,MUTED,False,JP)),
                          P(("財務数値は前提を明示した想定モデルであり、将来の達成を保証しません。",9.5,MUTED,False,JP))],line_spacing=1.2); set_ea(t)
footer(s)

out = "/home/user/regalis-hp/investor-materials/Regalis_投資家向け事業計画書_2026.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
